from dataclasses import dataclass
from typing import List, Dict, Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = 20.0
SLIDE_H = 11.25

BLUE = RGBColor(0x00, 0x2A, 0x6D)
CYAN = RGBColor(0x43, 0xCE, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x4A, 0x4F, 0x57)
GRID = RGBColor(0xD9, 0xE2, 0xEF)

LEFT_W = 6.59
RIGHT_TITLE_X = 8.03
RIGHT_BULLET_X = 6.46
RIGHT_W = 13.54
CONTENT_TOP = 1.58
CONTENT_BOTTOM = 10.95

TITLE_X = 0.15
TITLE_W = 6.33
TITLE_Y = 1.26
TITLE_H = 2.05

FORM_X = 0.15
FORM_W = 6.33
FORM_Y = 2.92
FORM_H = 2.05

TOOLS_X = 0.26
TOOLS_W = 5.31
TOOLS_Y = 4.82
TOOLS_H = 3.23

SOFT_X = 0.26
SOFT_W = 5.31
SOFT_Y = 8.13
SOFT_H = 2.32

COMP_TITLE_X = 7.94
COMP_TITLE_Y = 0.43
COMP_TITLE_W = 4.1
COMP_LINES_X = 7.93
COMP_LINES_Y = 0.92
COMP_LINES_W = 6.2
LOGO_X = 16.0
LOGO_Y = 0.22
LOGO_W = 2.65

PLUS_X = 7.70
PLUS_W = 0.27


@dataclass(frozen=True)
class LayoutProfile:
    name: str
    title_pt: float
    bullet_pt: float
    title_chars: int
    bullet_chars: int
    title_gap: float
    mission_gap: float
    plus_size: float


PROFILES = [
    LayoutProfile("normal", 15.0, 12.0, 68, 120, 0.18, 0.05, 0.24),
    LayoutProfile("compact", 14.0, 11.2, 74, 130, 0.14, 0.04, 0.22),
    LayoutProfile("dense", 13.2, 10.6, 80, 140, 0.12, 0.03, 0.20),
]


def _rgb(value: RGBColor) -> RGBColor:
    return value


def _wrap_lines(text: str, max_chars: int) -> List[str]:
    words = str(text or "").strip().split()
    if not words:
        return [""]
    lines: List[str] = []
    current: List[str] = []
    current_len = 0
    for word in words:
        if not current:
            current = [word]
            current_len = len(word)
            continue
        if current_len + 1 + len(word) <= max_chars:
            current.append(word)
            current_len += 1 + len(word)
        else:
            lines.append(" ".join(current))
            current = [word]
            current_len = len(word)
    if current:
        lines.append(" ".join(current))
    return lines or [""]


def _lines_height(line_count: int, font_pt: float, spacing: float = 1.15, padding: float = 0.02) -> float:
    return max(line_count, 1) * (font_pt * spacing / 72.0) + padding


def _mission_title(mission: Dict[str, Any]) -> str:
    company = str(mission.get("entreprise", "")).strip()
    role = str(mission.get("poste", "")).strip()
    duration = str(mission.get("duree", "")).strip().replace(" - ", " - ")
    parts = [p for p in [company, role, duration] if p]
    return " - ".join(parts)


def _mission_height(mission: Dict[str, Any], profile: LayoutProfile) -> float:
    title_lines = _wrap_lines(_mission_title(mission), profile.title_chars)
    title_h = _lines_height(len(title_lines), profile.title_pt, spacing=1.08, padding=0.01)
    bullets = mission.get("realisations", [])[:5]
    bullet_line_count = 0
    for bullet in bullets:
        bullet_line_count += len(_wrap_lines(f"• {bullet}", profile.bullet_chars))
    bullets_h = _lines_height(max(bullet_line_count, 1), profile.bullet_pt, spacing=1.16, padding=0.03)
    return title_h + profile.title_gap + bullets_h


def _chunk_height(chunk: List[Dict[str, Any]], profile: LayoutProfile) -> float:
    if not chunk:
        return 0.0
    total = sum(_mission_height(item, profile) for item in chunk)
    total += profile.mission_gap * max(len(chunk) - 1, 0)
    return total


def _available_height() -> float:
    return CONTENT_BOTTOM - CONTENT_TOP


def _best_profile_for_chunk(chunk: List[Dict[str, Any]]) -> LayoutProfile | None:
    for profile in PROFILES:
        if _chunk_height(chunk, profile) <= _available_height():
            return profile
    return None


def _plan_slides(missions: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    missions = missions[:9]
    if not missions:
        return [{"missions": [], "profile": PROFILES[0]}]

    best = None

    def score(plan: List[Dict[str, Any]]):
        blanks = [_available_height() - _chunk_height(item["missions"], item["profile"]) for item in plan]
        return (
            len(plan),
            sum(1 for item in plan if len(item["missions"]) == 1),
            round(sum(max(blank, 0) ** 2 for blank in blanks), 6),
            round(sum(max(blank, 0) for blank in blanks), 6),
        )

    def search(start: int, current: List[Dict[str, Any]]):
        nonlocal best
        if start >= len(missions):
            candidate = list(current)
            current_score = score(candidate)
            if best is None or current_score < best[0]:
                best = (current_score, candidate)
            return
        for size in range(1, 4):
            if size > 2:
                break
            end = start + size
            if end > len(missions):
                break
            chunk = missions[start:end]
            profile = _best_profile_for_chunk(chunk)
            if profile is None:
                continue
            current.append({"missions": chunk, "profile": profile})
            search(end, current)
            current.pop()

    search(0, [])
    if best:
        return best[1]
    return [{"missions": [mission], "profile": PROFILES[-1]} for mission in missions]


def _add_shape(slide, shape_type, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    return shape


def _textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.10)
    tf.margin_right = Inches(0.10)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    return tf


def _set_run(run, text, font_name, font_size, color, bold=False, italic=False):
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def _add_para(tf, text, font_name, font_size, color, bold=False, italic=False, level=0, align=PP_ALIGN.LEFT, space_after=0):
    p = tf.paragraphs[0] if len(tf.paragraphs) == 1 and not tf.paragraphs[0].text and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_after = Pt(space_after)
    r = p.add_run()
    _set_run(r, text, font_name, font_size, color, bold=bold, italic=italic)
    return p


def _render_plus_icon(slide):
    _add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0.52, 0.43, 0.52, 0.12, fill=CYAN)
    _add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0.72, 0.23, 0.12, 0.52, fill=CYAN)


def _render_logo(slide):
    tf = _textbox(slide, LOGO_X, LOGO_Y, LOGO_W, 0.52)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    for text, color in [("Value", BLUE), ("4", CYAN), ("Finance", BLUE)]:
        r = p.add_run()
        _set_run(r, text, "Montserrat", 26, color, bold=True)


def _render_sidebar(slide, data: Dict[str, Any], slide_idx: int, total: int):
    _add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0.0, 0.0, LEFT_W, SLIDE_H, fill=BLUE)
    _render_plus_icon(slide)

    first_name = str(data.get("prenom", "Consultant")).strip() or "Consultant"
    title = str(data.get("titre", "Consultant Finance")).strip() or "Consultant Finance"
    xp = str(data.get("annees_xp", "")).strip()

    tf_title = _textbox(slide, TITLE_X, TITLE_Y, TITLE_W, TITLE_H)
    _add_para(tf_title, f"{first_name} - {title} ({slide_idx}/{total})", "Montserrat", 27, WHITE, bold=True)
    xp_text = f"{xp} d'experience en {title}".strip() if xp else f"Experience en {title}"
    _add_para(tf_title, xp_text, "Montserrat", 19, CYAN, italic=True)

    _render_sidebar_section(slide, FORM_X, FORM_Y, FORM_W, FORM_H, "FORMATION", data.get("formation", [])[:3])
    _render_sidebar_section(slide, TOOLS_X, TOOLS_Y, TOOLS_W, TOOLS_H, "OUTILS", data.get("outils", [])[:8])
    _render_sidebar_section(slide, SOFT_X, SOFT_Y, SOFT_W, SOFT_H, "APTITUDES", data.get("aptitudes", [])[:4], leading_blank=True)


def _render_sidebar_section(slide, x: float, y: float, w: float, h: float, label: str, items: List[str], leading_blank: bool = False):
    tf = _textbox(slide, x, y, w, h)
    if leading_blank:
        _add_para(tf, "", "Aptos", 1, WHITE)
    _add_para(tf, label, "Montserrat", 18, WHITE, bold=True, space_after=4)
    _add_para(tf, "", "Aptos", 1, WHITE)
    for item in items:
        _add_para(tf, f"• {item}", "Aptos", 13, WHITE)


def _split_competences(items: List[str]) -> List[str]:
    if not items:
        return []
    if len(items) <= 2:
        return [" ".join(items)]
    if len(items) == 3:
        return [" ".join(items[:2]), items[2]]
    return [" ".join(items[:2]), " ".join(items[2:4])]


def _render_right_header(slide, data: Dict[str, Any]):
    tf = _textbox(slide, COMP_TITLE_X, COMP_TITLE_Y, COMP_TITLE_W, 0.40)
    _add_para(tf, "COMPETENCES CLÉS", "Montserrat", 18, BLUE, bold=True)

    lines = _split_competences(list(data.get("competences_cles", [])))
    tf_comp = _textbox(slide, COMP_LINES_X, COMP_LINES_Y, COMP_LINES_W, 0.60)
    for line in lines:
        _add_para(tf_comp, line, "Montserrat", 13, CYAN, bold=True)
    _render_logo(slide)


def _render_missions(slide, missions: List[Dict[str, Any]], profile: LayoutProfile):
    y = CONTENT_TOP
    for mission in missions:
        height = _mission_height(mission, profile)
        title_text = _mission_title(mission)
        title_lines = _wrap_lines(title_text, profile.title_chars)
        title_h = _lines_height(len(title_lines), profile.title_pt, spacing=1.08, padding=0.01)

        plus_tf = _textbox(slide, PLUS_X, y + 0.03, PLUS_W, 0.24)
        _add_para(plus_tf, "+", "Montserrat", 26 if profile.name == "normal" else 24, CYAN, bold=True)

        tf_title = _textbox(slide, RIGHT_TITLE_X, y, RIGHT_W - 2.7, title_h + 0.02)
        _add_para(tf_title, title_text, "Montserrat", profile.title_pt, BLUE, bold=True)

        bullets = mission.get("realisations", [])[:5]
        bullet_y = y + title_h + profile.title_gap
        bullet_h = max(height - title_h - profile.title_gap, 0.28)
        tf_bullets = _textbox(slide, RIGHT_BULLET_X, bullet_y, RIGHT_W, bullet_h)
        for bullet in bullets:
            _add_para(tf_bullets, f"• {bullet}", "Aptos", profile.bullet_pt, TEXT)

        y += height + profile.mission_gap


def _render_footer(slide, page_number: int):
    tf = _textbox(slide, 19.15, 10.55, 0.25, 0.18)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    _set_run(r, str(page_number), "Montserrat", 10, BLUE)
    slide.shapes.add_connector(1, Inches(19.45), Inches(10.64), Inches(19.78), Inches(10.64)).line.color.rgb = GRID


def _build_slide(slide, data: Dict[str, Any], slide_idx: int, total: int, missions: List[Dict[str, Any]], profile: LayoutProfile):
    _render_sidebar(slide, data, slide_idx, total)
    _render_right_header(slide, data)
    _render_missions(slide, missions, profile)
    _render_footer(slide, slide_idx)


def build_pro_pptx(data: Dict[str, Any], output_path):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    plan = _plan_slides(list(data.get("missions", [])))
    total = len(plan)
    blank = prs.slide_layouts[6]

    for idx, item in enumerate(plan, start=1):
        slide = prs.slides.add_slide(blank)
        _build_slide(slide, data, idx, total, item["missions"], item["profile"])

    prs.save(str(output_path))


def build_reference_template(output_path):
    sample = {
        "prenom": "Consultant",
        "titre": "Titre Consultant",
        "annees_xp": "5 ans",
        "competences_cles": ["Controle interne", "Analyse des risques", "Reporting financier"],
        "formation": ["Diplome - Ecole", "Master - Ecole", "Certification - Organisme"],
        "outils": ["Excel", "Power BI", "SAP", "Power Query"],
        "aptitudes": ["Rigueur", "Analyse", "Autonomie", "Esprit d'equipe"],
        "missions": [
            {
                "entreprise": "Entreprise",
                "poste": "Poste",
                "duree": "2023 - 2025",
                "realisations": [
                    "Exemple de realisation 1",
                    "Exemple de realisation 2",
                    "Exemple de realisation 3",
                ],
            }
        ],
    }
    build_pro_pptx(sample, output_path)
